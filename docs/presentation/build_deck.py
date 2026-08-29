#!/usr/bin/env python3
"""Generate the SYMBA T4.6 GA deck (corrected, code-verified content).

Approximates the SYMBA visual identity (white/blue/green, EU footer) — it is
NOT the official template; drop the content into the branded v9 template if
pixel-fidelity to the master is required. All figures verified against the
engine schemas (validation PASS) and the i18n UI strings.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

HERE = os.path.dirname(os.path.abspath(__file__))
SHOTS = os.path.join(HERE, "screenshots")

BLUE       = RGBColor(0x29, 0xAB, 0xE2)
DARKBLUE   = RGBColor(0x16, 0x4A, 0x6B)
TEXTBLUE   = RGBColor(0x1F, 0x5C, 0x83)
GREEN      = RGBColor(0x8D, 0xC6, 0x3F)
GREENFILL  = RGBColor(0xEC, 0xF5, 0xDC)
YELLOW     = RGBColor(0xC5, 0xD9, 0x2D)
GREY       = RGBColor(0x4A, 0x4A, 0x4A)
LIGHTBLUE  = RGBColor(0xE6, 0xF3, 0xFB)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height


def _set(tf, size, color, bold=False, align=PP_ALIGN.LEFT, font="Calibri"):
    for p in tf.paragraphs:
        p.alignment = align
        for r in p.runs:
            r.font.size = Pt(size); r.font.color.rgb = color
            r.font.bold = bold; r.font.name = font


def textbox(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    return tf


def para(tf, text, size, color, bold=False, bullet=False, space_after=6,
         first=False, align=PP_ALIGN.LEFT):
    p = tf.paragraphs[0] if first and not tf.paragraphs[0].runs else tf.add_paragraph()
    p.alignment = align; p.space_after = Pt(space_after)
    if bullet:
        r = p.add_run(); r.text = "—  "
        r.font.size = Pt(size); r.font.color.rgb = GREEN; r.font.bold = True
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.color.rgb = color; r.font.bold = bold
    r.font.name = "Calibri"
    return p


def banner(slide, title):
    bar = slide.shapes.add_shape(MSO_SHAPE.PENTAGON, Inches(0.0), Inches(0.18),
                                 Inches(11.0), Inches(0.85))
    bar.fill.solid(); bar.fill.fore_color.rgb = BLUE; bar.line.fill.background()
    tf = bar.text_frame; tf.margin_left = Inches(0.45)
    p = tf.paragraphs[0]; r = p.add_run(); r.text = title
    r.font.size = Pt(26); r.font.bold = True; r.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.LEFT; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    # SYMBA wordmark (text stand-in for logo)
    logo = textbox(slide, Inches(11.3), Inches(0.18), Inches(1.9), Inches(0.85),
                   MSO_ANCHOR.MIDDLE)
    p = logo.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = "SYM"; r.font.size = Pt(28); r.font.bold = True
    r.font.color.rgb = BLUE
    r2 = p.add_run(); r2.text = "BA"; r2.font.size = Pt(28); r2.font.bold = True
    r2.font.color.rgb = GREEN


def footer(slide):
    ln = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(6.78),
                                Inches(12.33), Pt(2.2))
    ln.fill.solid(); ln.fill.fore_color.rgb = YELLOW; ln.line.fill.background()
    tf = textbox(slide, Inches(0.5), Inches(6.9), Inches(9.5), Inches(0.5))
    para(tf, "Funded by the European Union — Horizon Europe R&I Programme · "
             "Grant Agreement N. 101135562", 9, GREY, first=True)
    tf2 = textbox(slide, Inches(10.2), Inches(6.9), Inches(2.9), Inches(0.5))
    para(tf2, "www.symbaproject.eu", 9, TEXTBLUE, bold=True, first=True,
         align=PP_ALIGN.RIGHT)


def green_box(slide, x, y, w, h):
    b = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    b.fill.solid(); b.fill.fore_color.rgb = GREENFILL
    b.line.color.rgb = GREEN; b.line.width = Pt(1.25)
    b.shadow.inherit = False
    return b


def blue_box(slide, x, y, w, h):
    b = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    b.fill.solid(); b.fill.fore_color.rgb = LIGHTBLUE
    b.line.color.rgb = BLUE; b.line.width = Pt(1.0); b.shadow.inherit = False
    return b


def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


# ===================================================================== SLIDE 1
s = prs.slides.add_slide(BLANK)
banner(s, "The problem: the methodology must be fixed first")
tf = textbox(s, Inches(0.5), Inches(1.25), Inches(7.4), Inches(0.9))
para(tf, "An LCSA of an industrial-symbiosis case (LCA + LCC + S-LCA) requires "
         "fixing the exact methodological pathway before any calculation.",
     16, DARKBLUE, bold=True, first=True)
tf = textbox(s, Inches(0.5), Inches(2.3), Inches(7.4), Inches(3.6))
para(tf, "Hundreds of interdependent choices", 15, TEXTBLUE, bold=True,
     bullet=True, first=True, space_after=2)
para(tf, "system boundary, allocation, zero-burden, futurisation, review "
         "level, pillar aggregation.", 13, GREY, space_after=12)
para(tf, "No single correct set-up", 15, TEXTBLUE, bold=True, bullet=True,
     space_after=2)
para(tf, "the choices constrain one another.", 13, GREY, space_after=12)
para(tf, "Expert-dependent and non-reproducible", 15, TEXTBLUE, bold=True,
     bullet=True, space_after=2)
para(tf, "two analysts reach two different set-ups for the same case.", 13,
     GREY)
green_box(s, Inches(8.25), Inches(1.4), Inches(4.55), Inches(3.0))
tf = textbox(s, Inches(8.55), Inches(1.6), Inches(3.95), Inches(2.6))
para(tf, "e.g.", 13, GREEN, bold=True, first=True, space_after=8)
para(tf, "A single fork can invert the result:", 14, DARKBLUE, bold=True,
     space_after=8)
para(tf, "the same residue treated as zero-burden waste (Q5=a) vs. as an "
         "economically-allocated co-product (Q5=c) — opposite impact "
         "attribution, opposite conclusion.", 13, GREY)
bb = blue_box(s, Inches(0.5), Inches(5.55), Inches(12.33), Inches(0.95))
tf = bb.text_frame; tf.word_wrap = True
para(tf, "Results that cannot be reproduced are not comparable — and not "
         "defensible.", 15, DARKBLUE, bold=True, first=True,
     align=PP_ALIGN.CENTER)
footer(s)
notes(s, "The hard part of an LCSA of an industrial-symbiosis case is not the "
         "computation — it is fixing the methodology before any number is "
         "produced. That set-up involves hundreds of interdependent decisions: "
         "system boundary, allocation, whether a waste enters with zero burden, "
         "whether parameters are projected to the future, what review level is "
         "required, whether pillars may be aggregated. These choices are not "
         "independent — one constrains the next — and there is no single correct "
         "answer. Today it depends on the individual analyst's judgement, so two "
         "competent experts set up the same case differently. The example on the "
         "right is the cleanest illustration: classify a residue as waste and you "
         "apply zero-burden; classify the very same residue as a co-product and "
         "you apply economic allocation — the environmental conclusion flips. A "
         "study you cannot reproduce is neither comparable nor defensible.")

# ===================================================================== SLIDE 2
s = prs.slides.add_slide(BLANK)
banner(s, "Method: case-level elicitation")
tf = textbox(s, Inches(0.5), Inches(1.2), Inches(12.3), Inches(0.55))
para(tf, "Seven case-level questions; each answer deterministically constrains "
         "a block of methodological choices.", 15, DARKBLUE, bold=True,
     first=True)
qs = [
    ("Q1", "What are you analyzing?  (required)",
     "specific exchange · eco-park/network · policy/programme · corporate report (ESG/CSRD) · monitoring"),
    ("Q2", "What phase is the system in?",
     "operational · under construction · design phase · baseline + N alternative scenarios"),
    ("Q3", "Which sustainability dimensions?  (required)",
     "ENV = LCA · ECO = LCC/MFCA/CBA/TEA · SOC = S-LCA"),
    ("Q4", "What is the report for?",
     "internal · external no-claim · public superiority claim · EU policy/PEF · academic"),
    ("Q5", "Nature of each symbiotic flow  (per flow)",
     "waste · free · co-product · interdependent · aggregated/black-box"),
    ("Q6", "Sector  +  Technology Readiness Level (TRL)",
     "14 canonical sectors + Other · TRL bands TRL9 … <5"),
    ("Q7", "Geographic spread",
     "co-located · regional · wide-area · multi-scale"),
]
y = 1.85
for tag, q, opts in qs:
    row = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5),
                             Inches(y), Inches(7.5), Inches(0.62))
    row.fill.solid(); row.fill.fore_color.rgb = WHITE
    row.line.color.rgb = BLUE; row.line.width = Pt(1.0); row.shadow.inherit = False
    tf = row.text_frame; tf.margin_left = Inches(0.15); tf.margin_top = Inches(0.02)
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = tag + "  "; r.font.bold = True; r.font.size = Pt(13)
    r.font.color.rgb = BLUE
    r = p.add_run(); r.text = q; r.font.bold = True; r.font.size = Pt(12.5)
    r.font.color.rgb = DARKBLUE
    p2 = tf.add_paragraph(); r = p2.add_run(); r.text = opts
    r.font.size = Pt(9.5); r.font.color.rgb = GREY
    y += 0.70
green_box(s, Inches(8.25), Inches(1.85), Inches(4.55), Inches(2.55))
tf = textbox(s, Inches(8.55), Inches(2.0), Inches(3.95), Inches(2.25))
para(tf, "e.g.", 13, GREEN, bold=True, first=True, space_after=6)
para(tf, "Q1 = Eco-park / network  →  ILCD Situation A (multi-actor)", 12.5,
     DARKBLUE, bold=True, space_after=8)
para(tf, "Q5 = waste (A pays B)  →  waste-paradigm / zero-burden rules", 12.5,
     DARKBLUE, bold=True, space_after=8)
para(tf, "Q5 = co-product (B pays A)  →  economic allocation + PEF Circular "
         "Footprint Formula", 12.5, DARKBLUE, bold=True)
bb = blue_box(s, Inches(8.25), Inches(4.55), Inches(4.55), Inches(1.05))
tf = bb.text_frame; tf.word_wrap = True
para(tf, "An open, branching choice set  →  one deterministic, reproducible "
         "configuration.", 13, DARKBLUE, bold=True, first=True,
     align=PP_ALIGN.CENTER)
footer(s)
notes(s, "The method inverts the problem. Instead of asking the analyst to "
         "master the full decision space, we ask seven case-level questions — "
         "the exact wording from the tool is on screen. Each answer is a "
         "complexity reducer: it does not just record a fact, it deterministically "
         "fixes a whole block of downstream methodological choices. Q1 sets the "
         "ILCD decision context; Q2 the temporal treatment; Q3 gates the three "
         "pillars (note ECO can be LCC, MFCA, CBA or TEA); Q4 the review and "
         "disclosure regime; Q5 is answered per exchanged flow and sets the "
         "allocation paradigm; Q6 covers sector overlay and technology readiness; "
         "Q7 the transport treatment. Worked example: answering 'eco-park / "
         "network' to Q1 fixes ILCD Situation A multi-actor on its own; and the "
         "same physical flow routes to opposite allocation rules depending on "
         "whether it is declared waste or co-product in Q5 — which is exactly the "
         "fork from the previous slide, now resolved deterministically.")

# ===================================================================== SLIDE 3
s = prs.slides.add_slide(BLANK)
banner(s, "Engine: schema-driven derivation")
tf = textbox(s, Inches(0.5), Inches(1.2), Inches(12.3), Inches(0.5))
para(tf, "Logic encoded in 5 validated, version-controlled schema files — "
         "not hard-coded.", 15, DARKBLUE, bold=True, first=True)
stages = [("L0", "derive\nbase states"), ("L1", "block forbidden\ncombinations"),
          ("Pathway", "assign routing\nIS-01…IS-05"), ("Activate", "switch on\nprovisions"),
          ("L2", "validate\nconsistency"), ("L3", "surface\ndecision points")]
x = 0.5; cw = 2.0; gap = 0.04
for tag, sub in stages:
    ch = s.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(x), Inches(2.05),
                            Inches(cw), Inches(1.15))
    ch.fill.solid(); ch.fill.fore_color.rgb = BLUE; ch.line.fill.background()
    ch.shadow.inherit = False
    tf = ch.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = tag; r.font.bold = True; r.font.size = Pt(14)
    r.font.color.rgb = WHITE
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
    r = p2.add_run(); r.text = sub; r.font.size = Pt(9.5); r.font.color.rgb = WHITE
    x += cw + gap
cards = [
    ("186", "atomic nodes", "LCA 59 · LCC 61 · S-LCA 66", BLUE),
    ("58", "cross-method rules", "4 blocks · 20 IR · 10 CIR · 5 FU · 7 B · 12 CDP", BLUE),
    ("PASS", "validation", "204 automated tests · 0 critical", GREEN),
]
x = 0.6
for big, mid, small, col in cards:
    cb = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(3.7),
                            Inches(3.9), Inches(2.0))
    cb.fill.solid(); cb.fill.fore_color.rgb = WHITE
    cb.line.color.rgb = BLUE; cb.line.width = Pt(1.0); cb.shadow.inherit = False
    tf = cb.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = big; r.font.bold = True; r.font.size = Pt(40)
    r.font.color.rgb = col
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
    r = p2.add_run(); r.text = mid; r.font.bold = True; r.font.size = Pt(15)
    r.font.color.rgb = DARKBLUE
    p3 = tf.add_paragraph(); p3.alignment = PP_ALIGN.CENTER
    r = p3.add_run(); r.text = small; r.font.size = Pt(10.5); r.font.color.rgb = GREY
    x += 4.05
tf = textbox(s, Inches(0.5), Inches(5.95), Inches(12.3), Inches(0.5))
para(tf, "Knowledge base derived from SYMBA deliverables D4.1 (LCA), D4.2 (LCC), "
         "D4.3 (S-LCA); auditable and updatable.", 12, GREY, first=True,
     align=PP_ALIGN.CENTER)
footer(s)
notes(s, "Behind the seven questions is a schema-driven engine: the methodology "
         "lives in five validated, version-controlled schema files, not in code, "
         "so a methodologist can audit and update the logic without touching the "
         "software. A run is a deterministic six-stage pipeline: L0 derives the "
         "base state of each pillar; L1 blocks forbidden combinations and stops "
         "immediately if one fires; the pathway stage assigns the terminal IS "
         "routing; activate switches on the applicable provisions; L2 validates "
         "cross-pillar consistency; L3 surfaces the critical decision points the "
         "analyst must resolve. The knowledge base carries 186 atomic nodes "
         "(59 LCA, 61 LCC, 66 S-LCA) and 58 cross-method rules, and the whole "
         "engine passes 204 automated tests with zero critical findings.")

# ===================================================================== SLIDE 4
s = prs.slides.add_slide(BLANK)
banner(s, "Output: configuration + consistency checks")
tf = textbox(s, Inches(0.5), Inches(1.2), Inches(12.3), Inches(0.5))
para(tf, "Per case: a full methodological configuration with provenance, plus "
         "automated checks.", 15, DARKBLUE, bold=True, first=True)
img = os.path.join(SHOTS, "03-result.png")
if os.path.exists(img):
    s.shapes.add_picture(img, Inches(0.5), Inches(1.85), height=Inches(4.55))
tf = textbox(s, Inches(6.5), Inches(1.95), Inches(6.3), Inches(2.6))
para(tf, "Assigned pathway", 14, TEXTBLUE, bold=True, bullet=True, first=True,
     space_after=2)
para(tf, "+ per-pillar configuration (ILCD situation, LCC type, S-LCA state).",
     12, GREY, space_after=10)
para(tf, "All activated provisions", 14, TEXTBLUE, bold=True, bullet=True,
     space_after=2)
para(tf, "each traceable to the triggering answer.", 12, GREY, space_after=10)
para(tf, "Blocks and cross-method warnings", 14, TEXTBLUE, bold=True,
     bullet=True, space_after=2)
para(tf, "forbidden combinations, inter-pillar inconsistencies, critical "
         "decision points.", 12, GREY)
green_box(s, Inches(6.5), Inches(4.75), Inches(6.3), Inches(1.65))
tf = textbox(s, Inches(6.75), Inches(4.9), Inches(5.85), Inches(1.4))
para(tf, "e.g.  Sokka et al. (2011) — six-actor eco-park, ENV-only", 12.5,
     DARKBLUE, bold=True, first=True, space_after=6)
para(tf, "pathway IS-01 (extended) · ILCD Situation A (multi-actor) · "
         "154/186 provisions activated · 0 blocks · 1 warning — B-05: "
         "geographic spread requires explicit transport in the LCA inventory "
         "and in LCC.", 12, GREY)
footer(s)
notes(s, "What the analyst receives is not a score but a complete, justified "
         "methodological configuration. The screenshot is a live run. On the "
         "left: the assigned pathway and the per-pillar configuration — ILCD "
         "situation, LCC type, S-LCA state — plus the full list of activated "
         "provisions, and critically every provision is traceable back to the "
         "answer that triggered it, so there are no silent assumptions. The "
         "engine also runs the checks: blocks for forbidden combinations, and "
         "cross-method warnings for inter-pillar inconsistencies. The example is "
         "Sokka 2011, a six-actor eco-park run environment-only: pathway IS-01 "
         "extended, ILCD Situation A multi-actor, 154 of 186 provisions "
         "activated, no blocks, and one warning telling the analyst that because "
         "the actors are geographically spread, transport must be modelled "
         "explicitly in both the LCA inventory and the LCC. Decision, "
         "justification, checks and documentation in a single pass.")

# ===================================================================== SLIDE 5
s = prs.slides.add_slide(BLANK)
banner(s, "Status & outlook")
gb = green_box(s, Inches(0.5), Inches(1.4), Inches(6.0), Inches(4.0))
gb.fill.fore_color.rgb = WHITE; gb.line.color.rgb = BLUE
tf = textbox(s, Inches(0.85), Inches(1.7), Inches(5.4), Inches(3.4))
para(tf, "Contribution", 17, BLUE, bold=True, first=True, space_after=12)
for h, d in [("Standardised set-up", "across partners"),
             ("Full provenance", "of every methodological choice"),
             ("Lower expertise barrier", "for sound studies")]:
    p = para(tf, h, 14, DARKBLUE, bold=True, bullet=True, space_after=2)
    para(tf, d, 12.5, GREY, space_after=12)
gb = green_box(s, Inches(6.83), Inches(1.4), Inches(6.0), Inches(4.0))
gb.fill.fore_color.rgb = WHITE; gb.line.color.rgb = BLUE
tf = textbox(s, Inches(7.18), Inches(1.7), Inches(5.4), Inches(3.4))
para(tf, "Status", 17, BLUE, bold=True, first=True, space_after=12)
for h, d in [("MVP complete", "engine validated (204 tests)"),
             ("Web application", "5 languages (EN · IT · FR · DE · ES)"),
             ("Aligned with D4.1–D4.3", "editable rule files"),
             ("Next", "application to SYMBA case studies")]:
    para(tf, h, 14, DARKBLUE, bold=True, bullet=True, space_after=2)
    para(tf, d, 12.5, GREY, space_after=10)
bb = blue_box(s, Inches(0.5), Inches(5.6), Inches(12.33), Inches(0.95))
tf = bb.text_frame; tf.word_wrap = True
para(tf, "Methodology selection: from expert-dependent and non-reproducible  "
         "→  deterministic, traceable, auditable.", 15, DARKBLUE, bold=True,
     first=True, align=PP_ALIGN.CENTER, space_after=2)
para(tf, "WP4 · Task 4.6 (ENCO) · SYMBA · Grant Agreement 101135562", 10.5,
     GREY, align=PP_ALIGN.CENTER)
footer(s)
notes(s, "To close: the contribution is to turn methodology selection from an "
         "expert-dependent, non-reproducible step into a standardised one, with "
         "full provenance for every choice and a lower expertise barrier so "
         "less-experienced analysts still produce methodologically sound studies. "
         "Status: the MVP is complete and the engine is validated by 204 tests; "
         "it is a web application in five languages; the rule files are aligned "
         "with deliverables D4.1 to D4.3 and remain editable. The next step is to "
         "apply it to the consortium's real case studies. Thank you — happy to "
         "take questions.")

# ===================================================================== SLIDE 6
s = prs.slides.add_slide(BLANK)
banner(s, "Backup (Q&A): pathways vs. configuration")
tf = textbox(s, Inches(0.5), Inches(1.2), Inches(12.3), Inches(0.55))
para(tf, "Five IS pathways are the outer routing layer — not the size of the "
         "methodological space.", 15, DARKBLUE, bold=True, first=True)
rows = [
    ("Layer", "What it fixes", "Cardinality", True),
    ("Routing", "IS pathway (Q1 × Q2)", "5  (IS-01…IS-05)", False),
    ("Pillar states", "ILCD situation × LCC type × S-LCA state", "5 × 4 × 2", False),
    ("Node activation", "which of 186 nodes switch on + field values", "not enumerable", False),
    ("Per-flow", "valuation type per exchanged flow", "5 ^ (n flows)", False),
]
y = 1.95
for c1, c2, c3, head in rows:
    for cx, cw_, txt, al in [(0.5, 3.0, c1, PP_ALIGN.LEFT),
                             (3.5, 3.6, c2, PP_ALIGN.LEFT),
                             (7.1, 2.8, c3, PP_ALIGN.LEFT)]:
        cell = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(cx), Inches(y),
                                  Inches(cw_), Inches(0.6))
        cell.fill.solid()
        cell.fill.fore_color.rgb = BLUE if head else WHITE
        cell.line.color.rgb = BLUE; cell.line.width = Pt(0.75)
        cell.shadow.inherit = False
        tf = cell.text_frame; tf.margin_left = Inches(0.12)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]; p.alignment = al
        r = p.add_run(); r.text = txt
        r.font.size = Pt(12.5 if head else 12)
        r.font.bold = head or (cx == 0.5)
        r.font.color.rgb = WHITE if head else DARKBLUE
    y += 0.62
green_box(s, Inches(0.5), Inches(5.15), Inches(9.4), Inches(1.35))
tf = textbox(s, Inches(0.8), Inches(5.3), Inches(8.9), Inches(1.1))
para(tf, "The configuration is a high-dimensional vector assembled node by "
         "node — which is why the engine is schema-driven, not a lookup table. "
         "Each study resolves 150+ provisions deterministically.", 13.5,
     DARKBLUE, bold=True, first=True)
footer(s)
notes(s, "Backup slide for the predictable Q&A challenge: 'only five pathways?'. "
         "The five IS pathways are merely the outer routing decision. The full "
         "methodological configuration is layered: routing (5) on top of pillar "
         "states (5 ILCD situations × 4 LCC types × 2 S-LCA states) on top of the "
         "activation pattern of 186 nodes with their field values, on top of a "
         "per-flow valuation choice (five types for each exchanged flow). The "
         "result is a high-dimensional vector assembled node by node — not "
         "tabulable, which is precisely why the engine is schema-driven rather "
         "than a lookup table. The honest framing of the size: the input space "
         "runs into the thousands even before per-flow choices, which make it "
         "effectively unbounded — but the operative point is that each study "
         "resolves 150-plus provisions deterministically and reproducibly.")

OUT = os.path.join(HERE, "SYMBA_T46_GA_corrected.pptx")
prs.save(OUT)
print("saved", OUT)
